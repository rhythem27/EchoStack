import os
import json
import uuid
import asyncio
import logging
from concurrent.futures import ThreadPoolExecutor
from kafka import KafkaConsumer
from docling.document_converter import DocumentConverter
from langchain_text_splitters import MarkdownHeaderTextSplitter
from sentence_transformers import SentenceTransformer
from langfuse.decorators import observe, langfuse_context

from backend.config import settings
from backend.db import init_db_pool, close_db_pool, get_db_pool

# Setup logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger("backend-worker")

# Define markdown headers to split on
HEADERS_TO_SPLIT = [
    ("#", "Header_1"),
    ("##", "Header_2"),
    ("###", "Header_3"),
    ("####", "Header_4"),
]

class IngestionWorker:
    def __init__(self):
        import torch
        self.loop = asyncio.get_running_loop()
        self.executor = ThreadPoolExecutor(max_workers=3)
        
        device = "cuda" if torch.cuda.is_available() else "cpu"
        logger.info(f"Initializing SentenceTransformer BAAI/bge-small-en-v1.5 on device: {device}...")
        self.embed_model = SentenceTransformer("BAAI/bge-small-en-v1.5", device=device)
        
        logger.info("Initializing IBM Docling Document Converter with fast pipeline options...")
        try:
            from docling.datamodel.pipeline_options import PdfPipelineOptions
            from docling.document_converter import DocumentConverter, PdfFormatOption
            from docling.datamodel.base_models import InputFormat

            pipeline_options = PdfPipelineOptions()
            pipeline_options.do_ocr = False
            pipeline_options.do_table_structure = False

            self.doc_converter = DocumentConverter(
                format_options={
                    InputFormat.PDF: PdfFormatOption(pipeline_options=pipeline_options)
                }
            )
        except Exception:
            self.doc_converter = DocumentConverter()
        
        self.consumer = None
        if getattr(settings, 'ENABLE_KAFKA', False):
            try:
                logger.info("Initializing Kafka Consumer...")
                self.consumer = KafkaConsumer(
                    settings.KAFKA_INGESTION_TOPIC,
                    bootstrap_servers=settings.KAFKA_BOOTSTRAP_SERVERS.split(","),
                    group_id="document_processors",
                    value_deserializer=lambda m: json.loads(m.decode("utf-8")),
                    auto_offset_reset="earliest",
                    enable_auto_commit=True
                )
            except Exception as e:
                logger.warning(f"Kafka Consumer initialization skipped: {e}. Polling database mode active.")
        else:
            logger.info("Kafka Consumer disabled. Database polling mode active.")
        
        self.markdown_splitter = MarkdownHeaderTextSplitter(
            headers_to_split_on=HEADERS_TO_SPLIT,
            strip_headers=False
        )

    async def start(self):
        logger.info("Worker started successfully and listening for ingestion events...")
        try:
            while True:
                # 1. Poll Kafka if active
                if self.consumer is not None:
                    try:
                        msg_pack = self.consumer.poll(timeout_ms=200)
                        for tp, messages in msg_pack.items():
                            for message in messages:
                                payload = message.value
                                logger.info(f"Received Kafka ingestion event: {payload}")
                                await self.process_event(payload)
                    except Exception as kafka_poll_err:
                        logger.warning(f"Kafka poll error: {kafka_poll_err}")

                # 2. Database polling fallback for PENDING documents
                try:
                    pool = await get_db_pool()
                    async with pool.acquire() as conn:
                        rows = await conn.fetch(
                            "SELECT id, user_id, file_name FROM documents WHERE status = 'PENDING' LIMIT 5"
                        )
                        for row in rows:
                            doc_id_str = str(row['id'])
                            user_id_str = str(row['user_id'])
                            file_name = row['file_name']

                            # Find file path in uploads directory
                            upload_dir = settings.UPLOAD_DIR
                            matched_file = None
                            if os.path.exists(upload_dir):
                                for f in os.listdir(upload_dir):
                                    if f.startswith(doc_id_str):
                                        matched_file = os.path.join(upload_dir, f)
                                        break

                            if matched_file and os.path.exists(matched_file):
                                logger.info(f"Picked up PENDING document from PostgreSQL: {file_name} ({doc_id_str})")
                                payload = {
                                    "doc_id": doc_id_str,
                                    "user_id": user_id_str,
                                    "file_path": matched_file,
                                    "file_name": file_name
                                }
                                await self.process_event(payload)
                            else:
                                logger.warning(f"File for pending document {doc_id_str} not found in {upload_dir}. Marking FAILED.")
                                await conn.execute("UPDATE documents SET status = 'FAILED' WHERE id = $1", row['id'])
                except Exception as db_poll_err:
                    logger.error(f"Database polling error in worker: {db_poll_err}")

                # Yield execution thread briefly
                await asyncio.sleep(0.5)
        except asyncio.CancelledError:
            logger.info("Worker cancelled, exiting loop...")
        finally:
            if self.consumer:
                try:
                    self.consumer.close()
                except Exception:
                    pass

    @observe(name="document-ingestion-worker", as_type="span")
    async def process_event(self, payload: dict):
        doc_id_str = payload.get("doc_id")
        user_id_str = payload.get("user_id")
        file_path = payload.get("file_path")
        file_name = payload.get("file_name")

        if not all([doc_id_str, user_id_str, file_path]):
            logger.error(f"Incomplete event payload ignored: {payload}")
            return

        langfuse_context.update_current_observation(
            name="document-ingestion-worker",
            input={"doc_id": doc_id_str, "file_name": file_name},
            metadata={"user_id": user_id_str, "file_path": file_path}
        )

        doc_uuid = uuid.UUID(doc_id_str)
        pool = await get_db_pool()

        # 1. Update status to PROCESSING
        logger.info(f"Setting status of document {doc_id_str} to PROCESSING...")
        async with pool.acquire() as conn:
            await conn.execute(
                "UPDATE documents SET status = 'PROCESSING' WHERE id = $1",
                doc_uuid
            )

        try:
            # 2. Extract layout-aware markdown via Docling with fallback for plain formats
            logger.info(f"Running Docling converter on file ({file_name}): {file_path}...")
            if not os.path.exists(file_path):
                raise FileNotFoundError(f"File not found at: {file_path}")

            ext = os.path.splitext(file_name)[1].lower()

            if ext in [".txt", ".csv", ".md"]:
                logger.info(f"Fast-path reading plain text document ({ext})...")
                with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                    markdown_text = f.read()
            elif ext == ".pdf":
                logger.info(f"Fast-path extracting PDF document ({file_name})...")
                markdown_text = None
                try:
                    import pypdf
                    reader = pypdf.PdfReader(file_path)
                    pages_text = []
                    for idx, page in enumerate(reader.pages):
                        t = page.extract_text()
                        if t and t.strip():
                            pages_text.append(f"## Page {idx+1}\n\n{t.strip()}")
                    if pages_text:
                        markdown_text = "\n\n".join(pages_text)
                except Exception as pypdf_err:
                    logger.warning(f"pypdf extraction fallback: {pypdf_err}")

                if not markdown_text:
                    logger.info("Falling back to Docling converter for PDF...")
                    conversion_result = await self.loop.run_in_executor(
                        self.executor,
                        self.doc_converter.convert,
                        file_path
                    )
                    markdown_text = conversion_result.document.export_to_markdown()
            elif ext == ".docx":
                logger.info(f"Fast-path extracting DOCX document ({file_name})...")
                try:
                    import docx
                    doc = docx.Document(file_path)
                    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
                    markdown_text = "\n\n".join(paragraphs)
                except Exception:
                    conversion_result = await self.loop.run_in_executor(
                        self.executor,
                        self.doc_converter.convert,
                        file_path
                    )
                    markdown_text = conversion_result.document.export_to_markdown()
            elif ext == ".pptx":
                logger.info(f"Fast-path extracting PPTX document ({file_name})...")
                try:
                    import pptx
                    prs = pptx.Presentation(file_path)
                    slides_text = []
                    for s_idx, slide in enumerate(prs.slides):
                        slide_words = [shape.text.strip() for shape in slide.shapes if hasattr(shape, "text") and shape.text.strip()]
                        if slide_words:
                            slides_text.append(f"## Slide {s_idx+1}\n\n" + "\n".join(slide_words))
                    markdown_text = "\n\n".join(slides_text)
                except Exception:
                    conversion_result = await self.loop.run_in_executor(
                        self.executor,
                        self.doc_converter.convert,
                        file_path
                    )
                    markdown_text = conversion_result.document.export_to_markdown()
            else:
                logger.info(f"Running IBM Docling layout analysis for document ({ext})...")
                conversion_result = await self.loop.run_in_executor(
                    self.executor,
                    self.doc_converter.convert,
                    file_path
                )
                markdown_text = conversion_result.document.export_to_markdown()
            
            # 3. Chunk using MarkdownHeaderTextSplitter
            logger.info("Splitting document content into markdown-header aware chunks...")
            chunks = self.markdown_splitter.split_text(markdown_text)
            
            if not chunks:
                logger.warning(f"No chunks extracted from document: {file_name}. Inserting full text as single chunk.")
                from langchain_core.documents import Document
                chunks = [Document(page_content=markdown_text)]

            # 4. Generate BGE small embeddings on GPU
            logger.info(f"Generating {len(chunks)} embeddings on GPU via BAAI/bge-small-en-v1.5...")
            texts = [chunk.page_content for chunk in chunks]
            
            embeddings = await self.loop.run_in_executor(
                self.executor,
                lambda: self.embed_model.encode(texts, convert_to_numpy=True).tolist()
            )

            # 5. Insert chunks, metadata & embeddings in PostgreSQL via asyncpg
            logger.info(f"Persisting vectors, metadata & chunks in Postgres vector_knowledge table...")
            async with pool.acquire() as conn:
                async with conn.transaction():
                    for idx, (chunk, embedding) in enumerate(zip(chunks, embeddings)):
                        vector_str = "[" + ",".join(map(str, embedding)) + "]"
                        
                        # Extract section header title if available from MarkdownHeaderTextSplitter
                        header_title = " > ".join([v for k, v in chunk.metadata.items() if k.startswith("Header_")]) if getattr(chunk, 'metadata', None) else "General"
                        
                        chunk_meta = {
                            "file_name": file_name,
                            "file_format": ext.lstrip("."),
                            "chunk_index": idx,
                            "section_title": header_title or "General",
                            "doc_id": doc_id_str
                        }


                        await conn.execute(
                            """
                            INSERT INTO vector_knowledge (doc_id, chunk_text, metadata, embedding)
                            VALUES ($1, $2, $3::jsonb, $4::vector)
                            """,
                            doc_uuid, chunk.page_content, json.dumps(chunk_meta), vector_str
                        )
                    
                    # 6. Update status to COMPLETE
                    await conn.execute(
                        "UPDATE documents SET status = 'COMPLETE' WHERE id = $1",
                        doc_uuid
                    )
            
            logger.info(f"Ingestion pipeline completed successfully for document {doc_id_str}.")
            langfuse_context.update_current_observation(output={"status": "COMPLETE", "chunks": len(chunks)})

        except Exception as e:
            logger.error(f"Ingestion processing failed for document {doc_id_str}: {e}")
            langfuse_context.update_current_observation(output={"status": "FAILED", "error": str(e)})
            async with pool.acquire() as conn:
                await conn.execute(
                    "UPDATE documents SET status = 'FAILED' WHERE id = $1",
                    doc_uuid
                )
        finally:
            if os.path.exists(file_path):
                try:
                    os.remove(file_path)
                    logger.info(f"Removed temporary upload file: {file_path}")
                except Exception as cleanup_err:
                    logger.error(f"Failed to remove temp file {file_path}: {cleanup_err}")

async def main():
    await init_db_pool()
    worker = IngestionWorker()
    try:
        await worker.start()
    finally:
        try:
            langfuse_context.flush()
        except Exception as e:
            logger.error(f"Failed to flush Langfuse context in worker: {e}")
        await close_db_pool()

if __name__ == "__main__":
    try:
        asyncio.run(main())
    except KeyboardInterrupt:
        logger.info("Worker shut down by keyboard interrupt.")
